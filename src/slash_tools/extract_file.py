from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from .context import SlashToolContext
from .errors import SlashToolError
from .path_safety import output_root, require_working_folder, resolve_inside
from .results import SlashToolResult


SUPPORTED_EXTRACT_EXTENSIONS = {".xlsx", ".pptx", ".pdf", ".docx"}


def extract_file_command(
    args: list[str],
    working_folder: str | Path | None,
    context: SlashToolContext,
    progress=None,
) -> SlashToolResult:
    if len(args) != 1:
        raise SlashToolError("사용법: /extract_file <xlsx|pptx|pdf|docx 경로>")
    output_path = extract_file_to_markdown(args[0], working_folder, context, progress=progress)
    return SlashToolResult(
        text=f"파일 추출 markdown이 저장되었습니다:\n{output_path}",
        tool_name="/extract_file",
    )


def extract_file_to_markdown(
    raw_path: str,
    working_folder: str | Path | None,
    context: SlashToolContext,
    progress=None,
) -> Path:
    root = require_working_folder(working_folder)
    source = resolve_inside(root, raw_path)
    if not source.exists():
        raise SlashToolError(f"파일이 존재하지 않습니다: {source}")
    if not source.is_file():
        raise SlashToolError(f"파일 경로가 아닙니다: {source}")
    if source.suffix.lower() not in SUPPORTED_EXTRACT_EXTENSIONS:
        raise SlashToolError(f"/extract_file에서 지원하지 않는 파일 형식입니다: {source.suffix}")

    if progress is not None:
        progress("status", f"{source.name} 파일을 추출하는 중...")
    context.check_cancelled()
    markdown = extract_file_markdown(source, context)
    out_dir = output_root(root, "extract_docs")
    output_path = out_dir / f"{source.name}.md"
    output_path.write_text(markdown, encoding="utf-8")
    return output_path


def extract_file_markdown(source: Path, context: SlashToolContext) -> str:
    extension = source.suffix.lower()
    if extension == ".xlsx":
        body = _extract_xlsx(source, context)
    elif extension == ".pptx":
        body = _extract_pptx(source, context)
    elif extension == ".pdf":
        body = _extract_pdf(source, context)
    elif extension == ".docx":
        body = _extract_docx(source, context)
    else:
        raise SlashToolError(f"/extract_file에서 지원하지 않는 파일 형식입니다: {source.suffix}")

    extracted_at = datetime.now(timezone.utc).isoformat()
    header = [
        f"# Extracted: {source.name}",
        "",
        f"- Source: `{source}`",
        f"- Type: `{extension.lstrip('.')}`",
        f"- Extracted At: `{extracted_at}`",
        "",
    ]
    return "\n".join(header) + body.strip() + "\n"


def _extract_xlsx(source: Path, context: SlashToolContext) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except Exception as exc:
        raise SlashToolError("XLSX 추출에는 openpyxl이 필요합니다.") from exc

    try:
        workbook = load_workbook(str(source), data_only=True, read_only=True)
    except Exception as exc:
        raise SlashToolError(f"XLSX 파일을 열지 못했습니다: {source.name}: {exc}") from exc

    sections: list[str] = []
    try:
        for sheet in workbook.worksheets:
            context.check_cancelled()
            rows: list[list[str]] = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value).replace("\n", " ").strip() for value in row]
                if any(values):
                    rows.append(values)
            sections.append(f"## Sheet: {sheet.title}\n\n{_markdown_table(rows)}")
    finally:
        workbook.close()
    return "\n\n".join(sections)


def _extract_pptx(source: Path, context: SlashToolContext) -> str:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except Exception as exc:
        raise SlashToolError("PPTX 추출에는 python-pptx가 필요합니다.") from exc

    try:
        presentation = Presentation(str(source))
    except Exception as exc:
        raise SlashToolError(f"PPTX 파일을 열지 못했습니다: {source.name}: {exc}") from exc

    sections: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        context.check_cancelled()
        lines = [f"## Slide {slide_index}", ""]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                lines.append(text.strip())
                lines.append("")
            if getattr(shape, "has_table", False):
                table_rows = []
                for row in shape.table.rows:
                    table_rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
                lines.append(_markdown_table(table_rows))
                lines.append("")
        sections.append("\n".join(lines).strip())
    return "\n\n".join(sections)


def _extract_pdf(source: Path, context: SlashToolContext) -> str:
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except Exception as exc:
        raise SlashToolError("PDF 추출에는 pypdf가 필요합니다.") from exc

    try:
        reader = PdfReader(str(source))
    except Exception as exc:
        raise SlashToolError(f"PDF 파일을 열지 못했습니다: {source.name}: {exc}") from exc

    sections: list[str] = []
    for page_index, page in enumerate(reader.pages, start=1):
        context.check_cancelled()
        text = page.extract_text() or ""
        sections.append(f"## Page {page_index}\n\n{text.strip()}")
    return "\n\n".join(sections)


def _extract_docx(source: Path, context: SlashToolContext) -> str:
    try:
        from docx import Document  # type: ignore[import-not-found]
    except Exception as exc:
        raise SlashToolError("DOCX 추출에는 python-docx가 필요합니다.") from exc

    try:
        document = Document(str(source))
    except Exception as exc:
        raise SlashToolError(f"DOCX 파일을 열지 못했습니다: {source.name}: {exc}") from exc

    blocks: list[str] = ["## Document Text", ""]
    for paragraph in document.paragraphs:
        context.check_cancelled()
        text = paragraph.text.strip()
        if text:
            blocks.append(text)
            blocks.append("")
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip().replace("\n", " ") for cell in row.cells])
        blocks.append(f"## Table {table_index}")
        blocks.append("")
        blocks.append(_markdown_table(rows))
        blocks.append("")
    return "\n".join(blocks)


def _markdown_table(rows: list[list[Any]]) -> str:
    if not rows:
        return "_No extractable tabular content._"
    width = max(len(row) for row in rows)
    normalized = [[str(cell) for cell in row] + [""] * (width - len(row)) for row in rows]
    header = normalized[0]
    body = normalized[1:]
    lines = [
        "| " + " | ".join(_escape_cell(cell) for cell in header) + " |",
        "| " + " | ".join("---" for _ in header) + " |",
    ]
    for row in body:
        lines.append("| " + " | ".join(_escape_cell(cell) for cell in row) + " |")
    return "\n".join(lines)


def _escape_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", " ").strip()
