from __future__ import annotations

<<<<<<< HEAD
from pathlib import Path
from typing import Any

from src.ingestion.parsers.base import DocumentParser, ParserError
from src.models.schemas import Asset, ParsedDocument, Section


class PptxParser(DocumentParser):
    """MVP PPTX parser that extracts slide text and image-like shape metadata."""

    file_type = "pptx"

    def parse(self, path: Path) -> ParsedDocument:
=======
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from src.ingestion.parsers.base import DocumentParser, ParserError
from src.ingestion.parsers.asset_utils import guess_mime_type, resolve_ooxml_target, write_asset_file
from src.models.schemas import Asset, ParsedDocument, Section

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PACKAGE_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


class PptxParser(DocumentParser):
    """PPTX parser that extracts slide text and embedded images."""

    file_type = "pptx"

    def parse(self, path: Path, asset_output_dir: Path | None = None) -> ParsedDocument:
>>>>>>> 4b1f4179239ca3b0466426fe629135dfeba590a3
        try:
            from pptx import Presentation  # type: ignore[import-not-found]
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]
        except Exception as exc:
            raise ParserError("PPTX parsing requires python-pptx.") from exc

        try:
            presentation = Presentation(str(path))
        except Exception as exc:
            raise ParserError(f"Failed to open PPTX {path}: {exc}") from exc

        doc_id = self.doc_id(path)
        sections: list[Section] = []
<<<<<<< HEAD
        assets: list[Asset] = []
=======
        assets = _extract_embedded_images(path, doc_id, asset_output_dir)
        assets_by_slide: dict[int, list[Asset]] = {}
        for asset in assets:
            if isinstance(asset.page_or_slide, int):
                assets_by_slide.setdefault(asset.page_or_slide, []).append(asset)
>>>>>>> 4b1f4179239ca3b0466426fe629135dfeba590a3
        slide_texts: list[str] = []

        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = [_shape_text(shape) for shape in slide.shapes]
            texts = [text for text in texts if text]
            title = _slide_title(slide, slide_number)
            slide_text = "\n".join(texts).strip()
            slide_texts.append(slide_text)

            section_assets: list[Asset] = []
            for shape_index, shape in enumerate(slide.shapes, start=1):
                if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
<<<<<<< HEAD
                    asset = Asset(
                        asset_id=f"{doc_id}-slide-{slide_number}-asset-{shape_index}",
                        kind="image",
                        source_file=str(path.resolve()),
                        page_or_slide=slide_number,
                        path="",
                        caption=getattr(shape, "name", ""),
                        metadata={
                            "shape_name": getattr(shape, "name", ""),
                            "width": int(getattr(shape, "width", 0)),
                            "height": int(getattr(shape, "height", 0)),
                        },
                    )
                    section_assets.append(asset)
                    assets.append(asset)
=======
                    matching = next(
                        (
                            asset
                            for asset in assets_by_slide.get(slide_number, [])
                            if asset.metadata.get("shape_name") == getattr(shape, "name", "")
                        ),
                        None,
                    )
                    if matching is not None and matching not in section_assets:
                        section_assets.append(matching)
            if not section_assets:
                section_assets = list(assets_by_slide.get(slide_number, []))
>>>>>>> 4b1f4179239ca3b0466426fe629135dfeba590a3

            sections.append(
                Section(
                    section_id=f"{doc_id}-slide-{slide_number}",
                    title=title,
                    level=1,
                    text=slide_text,
                    page_or_slide=slide_number,
                    assets=section_assets,
                    metadata={"parser": "pptx", "shape_count": len(slide.shapes)},
                )
            )

        return ParsedDocument(
            doc_id=doc_id,
            file_path=str(path.resolve()),
            file_type=self.file_type,
            title=self.title_from_path(path),
            text="\n\n".join(text for text in slide_texts if text).strip(),
            sections=sections,
            assets=assets,
            metadata={"slide_count": len(presentation.slides)},
        )


def _shape_text(shape: Any) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    return "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())


def _slide_title(slide: Any, slide_number: int) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is not None:
        title = _shape_text(title_shape)
        if title:
            return title.splitlines()[0]
    return f"Slide {slide_number}"

<<<<<<< HEAD
=======

def _extract_embedded_images(path: Path, doc_id: str, asset_output_dir: Path | None) -> list[Asset]:
    assets: list[Asset] = []
    try:
        archive = zipfile.ZipFile(path)
    except Exception as exc:
        raise ParserError(f"Failed to inspect embedded PPTX assets in {path}: {exc}") from exc

    with archive:
        media_parts = {
            name: archive.read(name)
            for name in archive.namelist()
            if name.startswith("ppt/media/") and not name.endswith("/")
        }
        for slide_number in range(1, 10000):
            slide_part = f"ppt/slides/slide{slide_number}.xml"
            if slide_part not in archive.namelist():
                break
            slide_assets = _extract_slide_assets(
                archive,
                media_parts,
                slide_part,
                slide_number,
                doc_id,
                path,
                asset_output_dir,
            )
            assets.extend(slide_assets)
    return assets


def _extract_slide_assets(
    archive: zipfile.ZipFile,
    media_parts: dict[str, bytes],
    slide_part: str,
    slide_number: int,
    doc_id: str,
    path: Path,
    asset_output_dir: Path | None,
) -> list[Asset]:
    rels_part = f"ppt/slides/_rels/slide{slide_number}.xml.rels"
    rels = _read_relationships(archive, rels_part, slide_part)
    root = ET.fromstring(archive.read(slide_part))
    assets: list[Asset] = []
    for image_index, picture in enumerate(root.findall(f".//{{{P_NS}}}pic"), start=1):
        blip = picture.find(f".//{{{A_NS}}}blip")
        if blip is None:
            continue
        embed = blip.attrib.get(f"{{{R_NS}}}embed", "")
        target = rels.get(embed, "")
        data = media_parts.get(target)
        if not data:
            continue
        doc_pr = picture.find(f".//{{{P_NS}}}cNvPr")
        ext = picture.find(f".//{{{A_NS}}}xfrm/{{{A_NS}}}ext")
        asset_id = f"{doc_id}-slide-{slide_number}-image-{image_index:04d}"
        source_name = Path(target).name or f"slide-{slide_number}-image-{image_index}.bin"
        stored_path, sha256 = write_asset_file(asset_output_dir, doc_id, asset_id, source_name, data)
        width = _optional_int(ext.attrib.get("cx") if ext is not None else None)
        height = _optional_int(ext.attrib.get("cy") if ext is not None else None)
        assets.append(
            Asset(
                asset_id=asset_id,
                kind="image",
                source_file=str(path.resolve()),
                page_or_slide=slide_number,
                path=stored_path,
                caption=(doc_pr.attrib.get("descr", "") if doc_pr is not None else "") or (
                    doc_pr.attrib.get("name", "") if doc_pr is not None else ""
                ),
                mime_type=guess_mime_type(source_name),
                sha256=sha256,
                width=width,
                height=height,
                metadata={
                    "shape_name": doc_pr.attrib.get("name", "") if doc_pr is not None else "",
                    "relationship_id": embed,
                    "package_path": target,
                },
            )
        )
    return assets


def _read_relationships(archive: zipfile.ZipFile, rels_part: str, base_part: str) -> dict[str, str]:
    if rels_part not in archive.namelist():
        return {}
    root = ET.fromstring(archive.read(rels_part))
    rels: dict[str, str] = {}
    for rel in root.findall(f".//{{{PACKAGE_NS}}}Relationship"):
        rel_id = rel.attrib.get("Id", "")
        target = rel.attrib.get("Target", "")
        if rel_id and target:
            rels[rel_id] = resolve_ooxml_target(base_part, target)
    return rels


def _optional_int(value: str | None) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None
>>>>>>> 4b1f4179239ca3b0466426fe629135dfeba590a3
