import os
import tempfile
import time
import unittest
from pathlib import Path
from unittest.mock import patch

os.environ.setdefault("QT_QPA_PLATFORM", "offscreen")

from PySide6.QtWidgets import QApplication

from src.gui.config import AppConfig
from src.gui.console_session import ConsoleConfig, ConsoleSessionError
from src.gui.database import ChatRepository
from src.gui.gui import MainWindow
from src.gui.token_handler import ModelPrompt


class FakeConsole:
    def __init__(self) -> None:
        self.config = ConsoleConfig("/bin/echo", "/tmp/model.gguf", ctx_size=2048)
        self.prompts: list[str | ModelPrompt] = []
        self.stopped = False

    def start(self) -> None:
        self.stopped = False

    def stop(self) -> None:
        self.stopped = True

    def stop_generation(self) -> None:
        self.stopped = True

    def ask(self, prompt: str | ModelPrompt) -> str:
        self.prompts.append(prompt)
        for _ in range(20):
            if self.stopped:
                self.stopped = False
                raise ConsoleSessionError("Generation stopped.")
            time.sleep(0.005)
        return "assistant answer"


def process_events(app: QApplication, cycles: int = 80) -> None:
    for _ in range(cycles):
        app.processEvents()
        time.sleep(0.005)


class GuiFlowTests(unittest.TestCase):
    def setUp(self) -> None:
        self.app = QApplication.instance() or QApplication([])

    def test_stop_then_next_send_works(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(llama_cli_path="/bin/echo", model_path="/tmp/model.gguf", backend="cli"),
        )

        window.input_edit.setPlainText("first")
        window.on_send()
        process_events(self.app, 5)
        window.on_stop_generation()
        process_events(self.app)

        self.assertIn("Generation stopped.", window.chat_view.toPlainText())
        self.assertTrue(window.send_btn.isEnabled())

        window.input_edit.setPlainText("second")
        window.on_send()
        process_events(self.app)

        self.assertIn("assistant answer", window.chat_view.toPlainText())
        prompt = console.prompts[-1]
        self.assertIsInstance(prompt, ModelPrompt)
        assert isinstance(prompt, ModelPrompt)
        self.assertEqual(prompt.messages[-1], {"role": "user", "content": "second"})
        self.assertIn("<start_of_turn>user\nsecond<end_of_turn>", prompt.completion_prompt)
        self.assertTrue(prompt.completion_prompt.endswith("<start_of_turn>model"))

        window.close()
        process_events(self.app, 5)

    def test_send_uses_recent_message_window(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        repository = ChatRepository(str(db_path))
        session_id = repository.create_session("Long chat")
        for index in range(6):
            repository.add_message(session_id, "user", f"user {index}")
            repository.add_message(session_id, "assistant", f"assistant {index}")

        window = MainWindow(
            console,
            repository,
            AppConfig(
                llama_cli_path="/bin/echo",
                model_path="/tmp/model.gguf",
                backend="cli",
                recent_message_limit=3,
            ),
        )
        window.current_session_id = session_id
        window.input_edit.setPlainText("current")

        window.on_send()
        process_events(self.app)

        prompt = console.prompts[-1]
        self.assertIsInstance(prompt, ModelPrompt)
        assert isinstance(prompt, ModelPrompt)
        prompt_text = "\n".join(message["content"] for message in prompt.messages)
        self.assertNotIn("user 0", prompt_text)
        self.assertNotIn("assistant 0", prompt_text)
        self.assertIn("assistant 5", prompt_text)
        self.assertIn("current", prompt_text)
        self.assertIn("Prompt context was shortened", window.chat_view.toPlainText())

        window.close()
        process_events(self.app, 5)

    def test_attach_folder_replaces_existing_attachment_list(self) -> None:
        console = FakeConsole()
        db_path = Path(tempfile.mkdtemp()) / "app.db"
        first_folder = Path(tempfile.mkdtemp())
        second_folder = Path(tempfile.mkdtemp())
        first_file = first_folder / "first.txt"
        second_file = second_folder / "second.pptx"
        first_file.write_text("first", encoding="utf-8")
        _write_sample_pptx(second_file)

        window = MainWindow(
            console,
            ChatRepository(str(db_path)),
            AppConfig(llama_cli_path="/bin/echo", model_path="/tmp/model.gguf", backend="cli"),
        )

        with patch("src.gui.gui.QFileDialog.getExistingDirectory", return_value=str(first_folder)):
            window.on_attach_files()
        self.assertEqual(window._attached_file_paths, [str(first_file.resolve())])

        with patch("src.gui.gui.QFileDialog.getExistingDirectory", return_value=str(second_folder)):
            window.on_attach_files()

        self.assertEqual(window._attached_file_paths, [str(second_file.resolve())])
        self.assertEqual(window.attachment_list.count(), 1)
        self.assertIn("second.pptx", window.attachment_list.item(0).text())

        window.close()
        process_events(self.app, 5)

def _write_sample_pptx(path: Path) -> None:
    try:
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise unittest.SkipTest("python-pptx is not installed") from exc

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "Deck"
    slide.placeholders[1].text = "Body"
    presentation.save(str(path))


if __name__ == "__main__":
    unittest.main()
